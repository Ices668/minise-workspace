from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLORS = {
    "navy": RGBColor(15, 23, 42),
    "blue": RGBColor(37, 99, 235),
    "cyan": RGBColor(8, 145, 178),
    "teal": RGBColor(13, 148, 136),
    "sky": RGBColor(224, 242, 254),
    "light": RGBColor(248, 250, 252),
    "muted": RGBColor(71, 85, 105),
    "text": RGBColor(30, 41, 59),
    "white": RGBColor(255, 255, 255),
    "accent": RGBColor(245, 158, 11),
    "green": RGBColor(22, 163, 74),
    "red": RGBColor(220, 38, 38),
}

slides = [
    {
        "title": "把 OpenClaw 用成个人工作台",
        "subtitle": "我的实战经验、工作流和技巧",
        "bullets": [
            "从\"会聊天\"到\"真能干活\"",
            "消息入口 + 工具编排 + 技能体系 + 工作流沉淀",
            "公司内部分享 / 2026-03-13",
        ],
        "tag": "OpenClaw 分享",
    },
    {
        "title": "为什么值得用 OpenClaw",
        "subtitle": "它解决的不是一个点问题，而是一整段工作链路的断裂",
        "bullets": [
            "工具太散：聊天、浏览器、终端、文档、文件夹来回切",
            "上下文断裂：需求、资料、产出不在一个地方，很难连续推进",
            "重复劳动高：同样的步骤每天都在重新解释、重新组织",
            "OpenClaw 的价值：把消息、工具、记忆、自动化串成一个工作台",
        ],
        "tag": "价值判断",
    },
    {
        "title": "我理解的 OpenClaw = 4 层能力栈",
        "subtitle": "不是一个聊天框，而是一套分层能力体系",
        "bullets": [
            "消息入口层：飞书等 IM 直接成为任务入口",
            "工具执行层：文件、Shell、Browser、Docs、设备能力可直接调用",
            "技能封装层：skill 把复杂能力变成可复用模块",
            "工作流沉淀层：把成功打法固定下来，后面越来越快",
        ],
        "tag": "能力模型",
    },
    {
        "title": "我现在的真实使用方式",
        "subtitle": "不是演示环境，是我日常就在跑的实际配置",
        "bullets": [
            "主通道：飞书直连，消息就是指令入口",
            "主工作区：minise-workspace，所有产出统一落盘",
            "主协作模式：对话驱动 + skill 调度 + 文件沉淀",
            "关键原则：能复用的都写进 memory、skill、脚本和规则里",
        ],
        "tag": "当前实践",
    },
    {
        "title": "我的核心工作流（总览）",
        "subtitle": "先判断任务类型，再走对应技能和交付路径",
        "bullets": [
            "需求进入：Boss 在飞书直接下达任务",
            "任务分流：信息检索 / 文档处理 / 设计展示 / 编码执行 / 自动化协作",
            "能力选择：优先 skill，其次工具，再到脚本",
            "产出沉淀：统一写入 workspace，必要时 push、回传、复盘",
        ],
        "tag": "流程视角",
    },
    {
        "title": "案例 1：产品设计展示工作流",
        "subtitle": "把需求、设计、确认、还原拆成阶段，避免返工",
        "bullets": [
            "enhance-prompt：把模糊需求打磨成可执行 prompt",
            "stitch-designer：先出设计稿，不急着写代码",
            "design-md：从设计稿抽设计系统，沉淀规范",
            "Boss 确认后，再进入 stitch-frontend 做工程还原",
        ],
        "tag": "案例 01",
    },
    {
        "title": "案例 2：工程还原与交付",
        "subtitle": "设计稿确认后再工程化，追求 1:1 还原",
        "bullets": [
            "stitch-frontend：按设计稿手写还原，避免自动化破坏细节",
            "shadcn-ui：做组件选型与集成，但不牺牲设计忠实度",
            "GitHub/push 脚本：工作区产出统一同步，方便追溯与复用",
            "铁律：图片保留、结构照抄、设计稿就是最终真相",
        ],
        "tag": "案例 02",
    },
    {
        "title": "案例 3：知识整理与复杂文档交付",
        "subtitle": "很多任务不是写系统，而是把复杂信息组织清楚",
        "bullets": [
            "用户旅程图 / Road Map 项目：资料理解 + 结构梳理 + 图形表达",
            "HTML 适合复杂排版展示，PDF 适合最终交付",
            "OpenClaw 擅长做：资料消化、内容重组、格式切换、持续迭代",
            "重点不是炫技，而是让信息表达更清楚、交付更快",
        ],
        "tag": "案例 03",
    },
    {
        "title": "从 CCG 视角看：为什么这套方式能持续进化",
        "subtitle": "能力不是一次性 prompt，而是可组合、可编排、可沉淀",
        "bullets": [
            "CCG 的价值在于：把单次回答升级成可持续执行系统",
            "skill 像插件，workflow 像 SOP，memory 像组织经验库",
            "一次做好一件事，不如把做成这件事的方法固定下来",
            "长期看，提效的关键不是更强模型，而是更稳的组织方式",
        ],
        "tag": "CCG 视角",
    },
    {
        "title": "我最常用的技能矩阵",
        "subtitle": "每个 skill 都应该有清晰边界，而不是一股脑全用",
        "bullets": [
            "enhance-prompt：把需求讲清楚",
            "stitch-designer / design-md：出设计、沉淀设计系统",
            "stitch-frontend / shadcn-ui：工程还原与组件集成",
            "github / skill-vetter / browser / playwright：同步、安全检查、自动化执行",
        ],
        "tag": "技能矩阵",
    },
    {
        "title": "真正有用的 8 个实战技巧",
        "subtitle": "这些不是概念，是我现在已经在用的方法",
        "bullets": [
            "1. 所有产出统一进 workspace，别让文件散落四处",
            "2. 先定工作流，再让模型干活",
            "3. 复杂任务拆阶段，别一口气让它全做完",
            "4. 把经验写进 MEMORY.md、skill 和脚本里",
            "5. 用 skill 代替重复解释",
            "6. 用消息通道驱动任务，而不是频繁切工具",
            "7. 对外动作谨慎，对内动作果断",
            "8. 每次交付后顺手把方法沉淀下来",
        ],
        "tag": "方法论",
    },
    {
        "title": "哪些场景最适合上 OpenClaw",
        "subtitle": "适合的场景越清楚，落地越顺",
        "bullets": [
            "频繁跨工具的知识工作",
            "需要沉淀个人或团队经验的方法型工作",
            "重复型任务，但每次又带一点变化",
            "强调边做边留痕、可追溯、可复盘的协作场景",
        ],
        "tag": "适用边界",
    },
    {
        "title": "哪些坑我踩过",
        "subtitle": "别把工具能力误以为就是完整工作流能力",
        "bullets": [
            "工具能用，不代表流程就顺；流程不顺，越强的工具越乱",
            "太早工程化，后面返工成本会很高",
            "没有统一产出目录，结果就是到处找文件",
            "截图、消息、附件格式要提前考虑，不然临门一脚翻车",
            "平台能力存在边界，替代方案要预先准备",
        ],
        "tag": "避坑经验",
    },
    {
        "title": "一句话总结",
        "subtitle": "OpenClaw 最强的不是回答问题，而是把做事的方法变成可执行系统",
        "bullets": [
            "从助手，到工作台，到个人作战系统",
            "真正的价值不是更像人，而是更像一个可靠的执行环境",
            "把经验固化下来，后面的每一次交付都会更轻松",
        ],
        "tag": "结论",
    },
    {
        "title": "Q&A / 可现场演示",
        "subtitle": "如果要演示，我建议挑一条最短链路当场跑通",
        "bullets": [
            "飞书发一句需求，直接进入工作台",
            "调用 skill 或工具执行一个真实动作",
            "产出落到 workspace，再回传结果",
            "最后展示 memory / skill / workflow 是如何形成闭环的",
        ],
        "tag": "结束页",
    },
]


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_block(slide, title, subtitle, tag):
    tag_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.45), Inches(1.8), Inches(0.45))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = COLORS['sky']
    tag_box.line.color.rgb = COLORS['sky']
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = tag
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = COLORS['cyan']
    p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(8.6), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = COLORS['navy']

    sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(2.02), Inches(10.5), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.size = Pt(15)
    r.font.color.rgb = COLORS['muted']


def add_bullets(slide, bullets):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.65), Inches(8.1), Inches(4.15))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS['white']
    panel.line.color.rgb = RGBColor(226, 232, 240)

    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.95), Inches(7.4), Inches(3.6))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22 if len(bullets) <= 4 else 18)
        p.font.color.rgb = COLORS['text']
        p.bullet = True
        p.space_after = Pt(10)


def add_side_card(slide, title, lines, color):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.1), Inches(1.3), Inches(3.55), Inches(5.5))
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.color.rgb = color

    tbox = slide.shapes.add_textbox(Inches(9.45), Inches(1.7), Inches(2.8), Inches(0.6))
    p = tbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = COLORS['white']

    box = slide.shapes.add_textbox(Inches(9.45), Inches(2.35), Inches(2.75), Inches(3.9))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['white']
        p.space_after = Pt(10)


def side_lines(index):
    mapping = {
        0: ("分享目标", ["讲清 OpenClaw 到底能做什么", "展示真实工作流，而非概念堆砌", "给同事一套可直接参考的落地方法"], COLORS['blue']),
        1: ("典型痛点", ["工具切换频繁", "上下文容易断", "重复劳动高", "经验难沉淀"], COLORS['teal']),
        2: ("四层关系", ["入口决定触发成本", "工具决定执行力", "技能决定复用率", "工作流决定稳定性"], COLORS['cyan']),
        3: ("我的原则", ["消息即入口", "文件要落盘", "经验要沉淀", "产出可追溯"], COLORS['blue']),
        4: ("流程关键词", ["识别任务类型", "匹配 skill", "统一产出", "回传与复盘"], COLORS['teal']),
        5: ("阶段好处", ["先确认视觉方向", "减少无效开发", "设计系统可复用"], COLORS['cyan']),
        6: ("交付原则", ["1:1 还原", "图片不能丢", "结构不能乱", "代码可维护"], COLORS['blue']),
        7: ("经验结论", ["表达清楚很重要", "格式切换是常态", "先交付再优化"], COLORS['teal']),
        8: ("CCG 启发", ["能力可组合", "流程可复制", "经验可积累", "系统会越用越顺"], COLORS['cyan']),
        9: ("常用 skills", ["enhance-prompt", "stitch-designer", "design-md", "stitch-frontend", "github", "skill-vetter"], COLORS['blue']),
        10: ("核心方法", ["统一工作区", "拆阶段", "写记忆", "少重复解释"], COLORS['teal']),
        11: ("适用判断", ["跨工具任务", "复杂知识工作", "需要沉淀经验", "需要可追溯"], COLORS['cyan']),
        12: ("避坑提醒", ["别过早工程化", "别忽略交付格式", "别让文件散落", "别高估平台边界"], COLORS['red']),
        13: ("一句话", ["把做事的方法", "变成一个", "可执行、可复用、可迭代", "的系统"], COLORS['accent']),
        14: ("演示建议", ["挑最短链路", "一步跑通", "让观众看到闭环"], COLORS['green']),
    }
    return mapping[index]

for idx, data in enumerate(slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, COLORS['light'])
    deco = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.28))
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLORS['navy']
    deco.line.color.rgb = COLORS['navy']
    add_title_block(slide, data['title'], data['subtitle'], data['tag'])
    add_bullets(slide, data['bullets'])
    title, lines, color = side_lines(idx)
    add_side_card(slide, title, lines, color)

out = '/home/lvyi/work/minise-workspace/presentations/OpenClaw使用经验与技巧分享-2026-03-13.pptx'
prs.save(out)
print(out)
